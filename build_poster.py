"""Rebuild the SPIN poster: swap the six matplotlib figures for the new poster-quality PNGs
(fit to their slots, no distortion), renumber captions 1-6 in reading order, and apply the
mentor's text styling (justified body, italic captions). Preserves the template (header bar,
logos, section bars, colours)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
import copy

EMU = 914400
SRC = "/tmp/spin.pptx"
OUT = "/sessions/amazing-jolly-maxwell/mnt/outputs/SPIN_Poster_v2.pptx"
FIGDIR = "/sessions/amazing-jolly-maxwell/mnt/outputs/poster_figs"

prs = Presentation(SRC)
slide = prs.slides[0]

# ---- map old figure pictures (by rounded position) -> new png + AR ----
from PIL import Image
def ar(png):
    im = Image.open(png); return im.width / im.height
NEW = {
    (3.23, 9.51): (f"{FIGDIR}/fig_schematic.png",),
    (12.59, 9.34): (f"{FIGDIR}/fig_montage.png",),
    (16.81, 9.34): (f"{FIGDIR}/fig_regional.png",),
    (12.47, 16.64): (f"{FIGDIR}/fig_paired.png",),
    (24.31, 5.35): (f"{FIGDIR}/fig_testretest.png",),
    (25.03, 12.05): (f"{FIGDIR}/fig_stabilization.png",),
}
# per-figure target boxes (left, top, maxW, maxH) in inches — chosen to fill the columns
BOX = {
    "fig_schematic.png":    (0.95, 9.6, 9.9, 3.2),
    "fig_montage.png":      (12.5, 9.2, 3.7, 4.9),
    "fig_regional.png":     (16.05, 9.2, 7.5, 4.9),
    "fig_paired.png":       (12.5, 16.3, 6.7, 5.9),
    "fig_testretest.png":   (24.15, 5.2, 6.05, 5.75),
    "fig_stabilization.png":(24.3, 11.95, 10.6, 3.55),
}

def fit(bx, by, bw, bh, a):
    """Largest (w,h) with aspect a=w/h fitting in box, centered."""
    w = bw; h = w / a
    if h > bh:
        h = bh; w = h * a
    l = bx + (bw - w) / 2
    t = by + (bh - h) / 2
    return l, t, w, h

# collect pictures to replace (can't modify while iterating shape tree cleanly)
to_replace = []
for sh in list(slide.shapes):
    if sh.shape_type == 13:
        key = (round(sh.left / EMU, 2), round(sh.top / EMU, 2))
        if key in NEW:
            to_replace.append((sh, NEW[key][0]))

for sh, png in to_replace:
    name = png.split("/")[-1]
    bx, by, bw, bh = BOX[name]
    l, t, w, h = fit(bx, by, bw, bh, ar(png))
    # drop old picture
    el = sh._element; el.getparent().remove(el)
    slide.shapes.add_picture(png, Emu(int(l*EMU)), Emu(int(t*EMU)),
                             Emu(int(w*EMU)), Emu(int(h*EMU)))

# ---- text edits: renumber captions + restyle ----
def set_text_keep_style(shape, new_text):
    """Replace text in the first run, keep its formatting; clear extra runs."""
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = new_text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)

RENUM = {
    "Figure 1. Spatial distribution of aperiodic exponent for rest and movie conditions. Error bars represent SEM.":
        "Figure 3. Xon 7-channel montage (left) and aperiodic exponent by scalp region for rest vs. movie (right). Error bars represent SEM.",
    "Figure 3. Paired aperiodic exponent estimates during quiet rest and movie watching (n = 19).":
        "Figure 4. Paired aperiodic exponent estimates during quiet rest and movie watching (n = 19).",
    "Figure 6. Test–retest agreement of aperiodic exponent estimates across two sessions for quiet rest and movie watching.":
        "Figure 5. Test–retest agreement of aperiodic exponent estimates across two sessions for quiet rest and movie watching.",
    "Figure 7. Recording duration required for individual aperiodic exponent estimates to stabilize.":
        "Figure 6. Recording duration required for individual aperiodic exponent estimates to stabilize.",
}
CAP_STARTS = ("Figure 1.", "Figure 2.", "Figure 3.", "Figure 4.", "Figure 5.", "Figure 6.")
BODY_JUSTIFY_CONTAINS = (
    "AD biomarkers do not fully", "Participants and device", "Signal processing",
    "Quiet rest and movie watching produced similar", "This pilot suggests that Xon",
    "Aperiodic exponent decreased from frontal", "Quiet rest showed stronger test",
    "Group-level reliability was achieved",
)

for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    txt = sh.text_frame.text.strip()
    if not txt:
        continue
    # renumber captions
    if txt in RENUM:
        set_text_keep_style(sh, RENUM[txt])
        txt = RENUM[txt]
    # italicise captions
    if txt.startswith(CAP_STARTS):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.italic = True
    # justify body paragraphs
    if any(txt.startswith(s) for s in BODY_JUSTIFY_CONTAINS):
        for p in sh.text_frame.paragraphs:
            if p.runs:
                p.alignment = PP_ALIGN.JUSTIFY

prs.save(OUT)

# ---- strip the leftover PowerPoint comment ("LK 1" marker) for a clean export ----
import zipfile, re, shutil, os as _os
tmp = OUT + ".tmp"
with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        name = item.filename
        if name.startswith("ppt/comments/") or name.endswith("commentAuthors.xml") \
           or "modernComment" in name or name.startswith("ppt/authors"):
            continue  # drop comment parts
        data = zin.read(name)
        if name == "[Content_Types].xml":
            txt = data.decode("utf-8")
            txt = re.sub(r"<Override[^>]*comment[^>]*/>", "", txt, flags=re.I)
            txt = re.sub(r"<Override[^>]*author[^>]*/>", "", txt, flags=re.I)
            data = txt.encode("utf-8")
        if name == "ppt/slides/_rels/slide1.xml.rels":
            txt = data.decode("utf-8")
            txt = re.sub(r"<Relationship[^>]*(comment|author)[^>]*/>", "", txt, flags=re.I)
            data = txt.encode("utf-8")
        zout.writestr(item, data)
shutil.move(tmp, OUT)
print("saved", OUT, "(comment stripped)")
